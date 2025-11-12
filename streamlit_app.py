import io
import math
import re
from dataclasses import dataclass

import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- optional (better layout) ---
HAVE_RNA = False
try:
    import RNA  # pip install ViennaRNA (or viennarna)
    HAVE_RNA = True
except Exception:
    HAVE_RNA = False

APP_TITLE = "RNA dot-bracket → PowerPoint (centered layout)"

# -------------------- parsing --------------------
def parse_rnafold_txt(payload: str):
    """Return (seq, db) from RNAfold-like text."""
    lines = [ln.strip() for ln in payload.splitlines() if ln.strip()]
    if len(lines) < 2:
        raise ValueError("Expected at least two lines (sequence + structure).")
    seq = re.sub(r"\s+", "", lines[0]).upper()
    if not re.fullmatch(r"[ACGU]+", seq):
        raise ValueError("First line must be RNA sequence (A/C/G/U).")
    db = None
    for ln in lines[1:]:
        tok = ln.split()[0]
        if re.fullmatch(r"[().]+", tok):
            db = tok
            break
    if db is None:
        raise ValueError("No valid dot-bracket line found.")
    if len(seq) != len(db):
        raise ValueError(f"Length mismatch: seq={len(seq)} vs db={len(db)}")
    return seq, db


def pairs_from_dotbracket(db: str):
    stack, pairs = [], []
    for i, ch in enumerate(db):
        if ch == "(":
            stack.append(i)
        elif ch == ")":
            if not stack:
                raise ValueError("Unbalanced dot-bracket: extra ')'")
            j = stack.pop()
            pairs.append((j, i))
    if stack:
        raise ValueError("Unbalanced dot-bracket: extra '('")
    return sorted(pairs)


# -------------------- coordinates --------------------
def get_naview_coords(db: str):
    """ViennaRNA Naview coordinates (RNAplot-like)."""
    try:
        RNA.cvar.rna_plot_type = 0
    except Exception:
        pass
    coords = RNA.get_xy_coordinates(db)
    xs = [coords.get(i).X for i in range(len(db))]
    ys = [coords.get(i).Y for i in range(len(db))]
    return xs, ys


def circular_coords(n: int, radius=120.0):
    xs, ys = [], []
    for k in range(n):
        th = 2 * math.pi * (k / n) - math.pi / 2
        xs.append(radius * math.cos(th))
        ys.append(radius * math.sin(th))
    return xs, ys


# -------------------- normalization --------------------
@dataclass
class LayoutParams:
    slide_w_in: float = 13.333
    slide_h_in: float = 7.5
    margin_in: float = 0.6
    node_diam_in: float = 0.26


def normalize_to_slide(xs, ys, lp: LayoutParams):
    xmin, xmax = min(xs), max(xs)
    ymin, ymax = min(ys), max(ys)
    w = xmax - xmin or 1.0
    h = ymax - ymin or 1.0
    avail_w = lp.slide_w_in - 2 * lp.margin_in
    avail_h = lp.slide_h_in - 2 * lp.margin_in
    scale = min(avail_w / w, avail_h / h)
    xs_scaled = [(x - xmin) * scale for x in xs]
    ys_scaled = [(y - ymin) * scale for y in ys]

    # Center in slide (both directions)
    x_mid = (max(xs_scaled) + min(xs_scaled)) / 2
    y_mid = (max(ys_scaled) + min(ys_scaled)) / 2
    slide_mid_x = lp.slide_w_in / 2
    slide_mid_y = lp.slide_h_in / 2

    xs_in = [((x - x_mid) + slide_mid_x) for x in xs_scaled]
    ys_in = [((y_mid - y) + slide_mid_y) for y in ys_scaled]  # Flip Y
    return xs_in, ys_in


# -------------------- drawing --------------------
NUC_COLORS = {
    "A": RGBColor(255, 99, 71),    # tomato
    "C": RGBColor(65, 105, 225),   # royal blue
    "G": RGBColor(34, 139, 34),    # forest green
    "U": RGBColor(238, 130, 238),  # violet
}


def add_nucleotide(slide, x_in, y_in, label: str, lp: LayoutParams):
    d = Inches(lp.node_diam_in)
    left = Inches(x_in) - d / 2
    top = Inches(y_in) - d / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(247, 247, 247)
    shp.line.width = Pt(1.0)
    shp.line.color.rgb = NUC_COLORS.get(label, RGBColor(90, 90, 90))
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(25, 25, 25)
    p.alignment = 1
    return shp


def add_conn(slide, x1, y1, x2, y2, w_pt=0.75, col=RGBColor(160, 160, 160)):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.width = Pt(w_pt)
    c.line.color.rgb = col
    return c


def build_ppt(seq: str, db: str, use_vienna: bool, lp: LayoutParams) -> bytes:
    # Coordinates
    if use_vienna and HAVE_RNA:
        xs, ys = get_naview_coords(db)
    else:
        xs, ys = circular_coords(len(seq), radius=120.0)
    xs_in, ys_in = normalize_to_slide(xs, ys, lp)

    prs = Presentation()
    prs.slide_width = Inches(lp.slide_w_in)
    prs.slide_height = Inches(lp.slide_h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Backbone
    for i in range(len(seq) - 1):
        add_conn(slide, xs_in[i], ys_in[i], xs_in[i + 1], ys_in[i + 1],
                 w_pt=0.75, col=RGBColor(205, 205, 205))
    # Base pairs
    for i, j in pairs_from_dotbracket(db):
        add_conn(slide, xs_in[i], ys_in[i], xs_in[j], ys_in[j],
                 w_pt=1.0, col=RGBColor(120, 120, 120))
    # Nucleotides
    for i, base in enumerate(seq):
        add_nucleotide(slide, xs_in[i], ys_in[i], base, lp)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()


# -------------------- Streamlit UI --------------------
st.set_page_config(page_title=APP_TITLE, page_icon="🧬", layout="centered")
st.title(APP_TITLE)
st.caption(
    "Upload an RNAfold-style .txt file (first line: sequence, next: dot-bracket). "
    "Generates a PowerPoint with editable nucleotides, centered on the slide."
)

uploaded = st.file_uploader("Upload RNAfold text (.txt)", type=["txt"])
with st.expander("Options"):
    node_d = st.slider("Circle diameter (inches)", 0.16, 0.45, 0.26, 0.01)
    use_vienna = st.checkbox(
        f"Use ViennaRNA (Naview) coordinates (detected: {'Yes' if HAVE_RNA else 'No'})",
        value=HAVE_RNA,
        help="If installed, uses Naview layout identical to RNAplot."
    )

if st.button("Generate PowerPoint"):
    if not uploaded:
        st.error("Please upload a file.")
    else:
        try:
            text = uploaded.read().decode("utf-8")
            seq, db = parse_rnafold_txt(text)
            lp = LayoutParams(node_diam_in=node_d)
            pptx_bytes = build_ppt(seq, db, use_vienna, lp)
            st.success("✅ Done! The RNA structure is centered on the slide.")
            st.download_button(
                "Download PowerPoint (.pptx)",
                pptx_bytes,
                file_name="rna_structure_centered.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception as e:
            st.error(f"Error: {e}")
